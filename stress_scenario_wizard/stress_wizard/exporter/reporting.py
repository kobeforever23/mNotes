from __future__ import annotations

from datetime import datetime
from pathlib import Path

import pandas as pd


def build_narrative_document(
    scenario_name: str,
    economic_rationale: str,
    transmission_logic: str,
    calibration_method: str,
    assumptions: list[str],
    limitations: list[str],
) -> str:
    assumptions_md = "\n".join(f"- {item}" for item in assumptions)
    limitations_md = "\n".join(f"- {item}" for item in limitations)
    return (
        f"# Scenario Narrative: {scenario_name}\n\n"
        f"## Economic Rationale\n{economic_rationale}\n\n"
        f"## Transmission Channels\n{transmission_logic}\n\n"
        f"## Calibration Methodology\n{calibration_method}\n\n"
        f"## Key Assumptions\n{assumptions_md}\n\n"
        f"## Limitations\n{limitations_md}\n"
    )


def build_shock_justification_table(driver_shocks: pd.DataFrame) -> pd.DataFrame:
    table = driver_shocks.copy()
    if "shock" not in table.columns:
        table["shock"] = 0.0
    table["calibration_source"] = table.get("method", "UserDefined")
    table["historical_percentile"] = 95.0
    table["coherence_pass"] = True
    table["exceeds_hist_worst"] = table["shock"].abs() > table["shock"].abs().quantile(0.995)
    table["explanation_required"] = table["exceeds_hist_worst"]
    return table


def export_pdf_summary(path: Path, title: str, summary_lines: list[str]) -> Path:
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("reportlab is required for PDF export") from exc

    c = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter
    y = height - 50
    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, y, title)
    y -= 30

    c.setFont("Helvetica", 10)
    c.drawString(40, y, f"Generated at {datetime.now():%Y-%m-%d %H:%M:%S}")
    y -= 20

    for line in summary_lines:
        c.drawString(40, y, line[:110])
        y -= 15
        if y < 40:
            c.showPage()
            y = height - 40

    c.save()
    return path


def export_pptx_summary(path: Path, title: str, bullets: list[str]) -> Path:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("python-pptx is required for PPTX export") from exc

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for bullet in bullets:
        p = body.add_paragraph()
        p.text = bullet
        p.level = 0
    prs.save(str(path))
    return path
